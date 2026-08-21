from pathlib import Path
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_THEME_COLOR


ROOT = Path(r"E:\Hackathon\Build with AI")
SOURCE_IMAGE = Path(r"E:\IOT project\simulation image sketch.png")
OUT = ROOT / "LUNA_IEEE_Multimodal_Voice_Assistant_Deck.pptx"
ASSET_DIR = ROOT / "luna_deck_assets"

W, H = 13.333, 7.5
BG = "0B1220"
PANEL = "111C2E"
PANEL_2 = "16243A"
TEXT = "F4F7FB"
MUTED = "A7B3C5"
TEAL = "2DD4BF"
PURPLE = "8B5CF6"
AMBER = "F59E0B"
RED = "FB7185"
BLUE = "60A5FA"
GREEN = "4ADE80"
WHITE = "FFFFFF"


def rgb(hexv):
    return RGBColor.from_string(hexv)


def add_shape(slide, kind, x, y, w, h, fill=None, line=None, radius=False):
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid(); shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line); shape.line.width = Pt(1)
    return shape


def textbox(slide, text, x, y, w, h, size=18, color=TEXT, bold=False,
            font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
            margin=0.05, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(margin); tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin); tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.name = font; p.font.size = Pt(size); p.font.bold = bold; p.font.italic = italic; p.font.color.rgb = rgb(color)
    return box


def rich_text(slide, lines, x, y, w, h, size=16, color=TEXT, bullet=True, spacing=1.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.03)
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            label, body, c = item
        else:
            label, body, c = "", item, color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " if bullet else "") + (label + " " if label else "") + body
        p.font.name = "Aptos"; p.font.size = Pt(size); p.font.color.rgb = rgb(c)
        p.space_after = Pt(7); p.line_spacing = spacing
    return box


def add_bg(slide, section, number, title, subtitle=None):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(BG)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, W, 0.12, TEAL)
    textbox(slide, section.upper(), 0.55, 0.28, 3.2, 0.22, 8.5, TEAL, True)
    textbox(slide, f"{number:02d}", 12.1, 0.25, 0.65, 0.25, 9, MUTED, True, align=PP_ALIGN.RIGHT)
    textbox(slide, title, 0.55, 0.62, 12.1, 0.55, 25, TEXT, True)
    if subtitle:
        textbox(slide, subtitle, 0.58, 1.18, 11.5, 0.34, 11, MUTED)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.55, 7.16, 12.2, 0.012, "26364D")
    textbox(slide, "LUNA • ESP32-S3 multimodal voice assistant", 0.55, 7.22, 6.0, 0.15, 7.5, MUTED)
    textbox(slide, "Research presentation", 10.0, 7.22, 2.75, 0.15, 7.5, MUTED, align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, title, body, accent=TEAL, body_size=12.5):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, PANEL, "26364D")
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.06, h, accent)
    textbox(slide, title, x+0.18, y+0.14, w-0.3, 0.28, 12.5, TEXT, True)
    textbox(slide, body, x+0.18, y+0.52, w-0.3, h-0.62, body_size, MUTED)


def pill(slide, text, x, y, w, color=TEAL):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 0.32, color)
    textbox(slide, text, x, y+0.03, w, 0.2, 9.5, BG, True, align=PP_ALIGN.CENTER)


def crop_assets():
    ASSET_DIR.mkdir(exist_ok=True)
    im = Image.open(SOURCE_IMAGE).convert("RGB")
    # Crops are intentionally tied to the supplied sketch, so component visuals remain faithful to the user’s prototype.
    crops = {
        "esp32": (475, 315, 790, 920),
        "mic": (280, 120, 480, 360),
        "oled": (480, 80, 810, 320),
        "amp": (860, 100, 1100, 330),
        "speaker": (1110, 100, 1480, 420),
        "encoder": (150, 330, 430, 600),
        "touch": (850, 330, 1090, 610),
        "power": (380, 650, 890, 910),
        "controls": (1080, 470, 1460, 850),
        "led": (130, 600, 410, 850),
    }
    paths = {}
    for name, box in crops.items():
        crop = im.crop(box)
        crop = ImageOps.contain(crop, (900, 620), method=Image.Resampling.LANCZOS)
        canvas = Image.new("RGB", (900, 620), (17, 28, 46))
        canvas.paste(crop, ((900-crop.width)//2, (620-crop.height)//2))
        p = ASSET_DIR / f"{name}.png"; canvas.save(p); paths[name] = p
    full = ASSET_DIR / "sketch_full.png"; im.save(full); paths["full"] = full
    return paths


def picture(slide, path, x, y, w, h, line="26364D"):
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    if line:
        sh = slide.shapes[-1]
        sh.line.color.rgb = rgb(line); sh.line.width = Pt(1)


def add_arrow(slide, x1, y1, x2, y2, color=TEAL, width=2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(color); line.line.width = Pt(width)
    line.line.end_arrowhead = True


def make_deck():
    assets = crop_assets()
    prs = Presentation(); prs.slide_width = Inches(W); prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    # 1. Title
    s = prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb = rgb(BG)
    add_shape(s, MSO_SHAPE.RECTANGLE, 0, 0, W, 0.14, TEAL)
    add_shape(s, MSO_SHAPE.OVAL, 9.45, -1.0, 5.1, 5.1, "12253A")
    add_shape(s, MSO_SHAPE.OVAL, 10.55, -0.15, 3.2, 3.2, "1B3051")
    textbox(s, "LUNA", 0.75, 0.72, 4.0, 0.75, 40, TEAL, True)
    textbox(s, "A multimodal AI voice assistant\nwith touch-driven affective interaction", 0.78, 1.55, 7.7, 1.25, 28, TEXT, True)
    textbox(s, "ESP32-S3 embedded platform • local AI server • Wokwi functional simulation", 0.82, 3.05, 7.4, 0.38, 13, MUTED)
    pill(s, "RESEARCH PRESENTATION", 0.82, 3.7, 1.85, PURPLE)
    textbox(s, "Prashant Gupta  •  Amrit Kumar Sah\nChennai Institute of Technology, India", 0.82, 4.35, 4.8, 0.62, 12, TEXT)
    picture(s, assets["oled"], 8.15, 1.1, 4.4, 2.95, None)
    picture(s, assets["touch"], 8.95, 4.25, 2.65, 1.75, None)
    textbox(s, "01", 12.15, 7.15, 0.55, 0.18, 8, MUTED, True, align=PP_ALIGN.RIGHT)

    # 2. Motivation / gap
    s = prs.slides.add_slide(blank); add_bg(s, "Context", 2, "Why LUNA?", "Voice is powerful, but voice-only interaction is not enough for a compact embedded assistant.")
    card(s, 0.65, 1.75, 3.75, 2.0, "Voice-only friction", "Noise, accents, ambiguity, accidental activation, privacy concerns, and waiting on the full AI pipeline.", RED)
    card(s, 4.8, 1.75, 3.75, 2.0, "Research gap", "Existing work is fragmented across voice UX, multimodality, privacy, and local processing; few systems combine these on a small embedded device.", AMBER)
    card(s, 8.95, 1.75, 3.75, 2.0, "Design response", "Treat voice, touch, display, buttons, and encoder as co-equal inputs that converge on one interaction manager.", TEAL)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.75, 4.35, 11.95, 1.55, PANEL_2, "2B4261")
    textbox(s, "Core idea", 1.0, 4.65, 1.3, 0.25, 11, TEAL, True)
    textbox(s, "A tap should be able to produce an immediate, visible, audible response — without speech recognition, language understanding, or network access.", 1.0, 5.0, 10.9, 0.55, 20, TEXT, True)
    textbox(s, "Paper scope: system-level integration and interaction design, not a new ASR or LLM algorithm.", 1.0, 5.65, 10.8, 0.22, 10.5, MUTED, italic=True)

    # 3. Contributions
    s = prs.slides.add_slide(blank); add_bg(s, "Contribution", 3, "What the paper contributes", "Five engineering contributions anchor the prototype and make the novelty testable.")
    contribs = [
        ("01", "Compact architecture", "ESP32-S3 coordinates deterministic I/O while heavier AI services run on a local host.", TEAL),
        ("02", "Affective touch", "TTP223 events map to explicit happy / sad response states.", PURPLE),
        ("03", "Unified state model", "Voice, touch, buttons, encoder, display, RGB, and audio remain synchronized.", AMBER),
        ("04", "Modular AI pipeline", "faster-whisper → Qwen3/Ollama → Piper or Azure TTS.", BLUE),
        ("05", "Honest validation", "Simulation-verified behavior is separated from physical acoustic and latency measurements.", GREEN),
    ]
    for i, (num, title, body, c) in enumerate(contribs):
        x = 0.75 + (i % 3) * 4.15; y = 1.85 + (i // 3) * 2.05
        add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 3.65, 1.6, PANEL, "26364D")
        textbox(s, num, x+0.2, y+0.18, 0.48, 0.35, 18, c, True)
        textbox(s, title, x+0.82, y+0.18, 2.55, 0.28, 13.5, TEXT, True)
        textbox(s, body, x+0.2, y+0.67, 3.15, 0.68, 11.3, MUTED)
    textbox(s, "Novelty is in the integration and interaction design: touch becomes a semantic event, not just a hardware interrupt.", 0.85, 6.28, 11.4, 0.35, 14, TEAL, True, align=PP_ALIGN.CENTER)

    # 4. Architecture
    s = prs.slides.add_slide(blank); add_bg(s, "System", 4, "LUNA system architecture", "The embedded edge stays deterministic; the local server supplies replaceable AI services.")
    # left edge
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.8, 1.75, 3.25, 4.55, PANEL, TEAL)
    textbox(s, "ESP32-S3 EDGE DEVICE", 1.05, 1.98, 2.7, 0.3, 13, TEAL, True, align=PP_ALIGN.CENTER)
    for i, t in enumerate(["INMP441 microphone", "SH1106 OLED face", "TTP223 touch ×2", "KY-040 encoder", "Buttons + RGB LED", "I2S audio output"]):
        add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.1, 2.55+i*0.53, 2.65, 0.34, PANEL_2, "2B4261")
        textbox(s, t, 1.17, 2.61+i*0.53, 2.5, 0.19, 10.5, TEXT, align=PP_ALIGN.CENTER)
    # middle manager
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 4.8, 2.45, 3.0, 2.2, "182B46", PURPLE)
    textbox(s, "INTERACTION MANAGER", 5.0, 2.75, 2.6, 0.3, 14, PURPLE, True, align=PP_ALIGN.CENTER)
    textbox(s, "Event-driven state machine\n\nIDLE • LISTENING • THINKING\nHAPPY • SAD • SLEEP", 5.1, 3.25, 2.4, 0.95, 12.2, TEXT, True, align=PP_ALIGN.CENTER)
    # right server
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 8.55, 1.75, 3.85, 4.55, PANEL, BLUE)
    textbox(s, "LOCAL AI SERVER", 8.85, 1.98, 3.25, 0.3, 13, BLUE, True, align=PP_ALIGN.CENTER)
    services = [("STT", "faster-whisper", TEAL), ("LLM", "Qwen3 via Ollama", PURPLE), ("TTS", "Piper / Azure", AMBER), ("Bridge", "Python session API", GREEN)]
    for i, (a,b,c) in enumerate(services):
        y=2.65+i*0.75
        add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 9.0, y, 2.95, 0.5, PANEL_2, "2B4261")
        textbox(s, a, 9.17, y+0.12, 0.55, 0.2, 10.5, c, True)
        textbox(s, b, 9.8, y+0.12, 1.95, 0.2, 10.5, TEXT)
    add_arrow(s, 4.05, 3.95, 4.8, 3.55, TEAL, 2.5); add_arrow(s, 7.8, 3.55, 8.55, 3.55, BLUE, 2.5)
    textbox(s, "local network", 7.86, 3.12, 0.75, 0.22, 8.5, MUTED, align=PP_ALIGN.CENTER)
    textbox(s, "Touch reactions can complete here", 4.55, 5.05, 3.55, 0.3, 10.5, AMBER, True, align=PP_ALIGN.CENTER)

    # 5. hardware map
    s = prs.slides.add_slide(blank); add_bg(s, "Hardware", 5, "Hardware implementation and wiring", "The supplied Wokwi sketch maps every interaction path around the ESP32-S3 controller.")
    picture(s, assets["full"], 0.6, 1.55, 8.15, 5.35)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 9.1, 1.65, 3.45, 4.95, PANEL, "26364D")
    textbox(s, "Key hardware blocks", 9.38, 1.95, 2.9, 0.28, 14, TEAL, True)
    rich_text(s, [
        ("Compute", "ESP32-S3-WROOM-1 N8R8", TEAL),
        ("Input", "INMP441 I2S mic + TTP223 touch", PURPLE),
        ("Output", "SH1106 OLED + MAX98357A + 3W speaker", AMBER),
        ("Control", "KY-040 encoder + sleep / mode / chat buttons", BLUE),
        ("Power", "3.7 V LiPo + switch + 100k/100k monitor", GREEN),
    ], 9.38, 2.45, 2.78, 2.65, 10.8)
    textbox(s, "The simulation validates GPIO and event logic; it does not reproduce the full acoustic or capacitive electrical behavior.", 9.38, 5.5, 2.8, 0.7, 10.2, MUTED, italic=True)

    # 6. components
    s = prs.slides.add_slide(blank); add_bg(s, "Hardware", 6, "Components selected for a student-scale prototype", "Each part serves a distinct interaction or feedback role in the multimodal loop.")
    comps = [("INMP441", "Digital I2S voice capture", "mic", TEAL), ("SH1106 OLED", "Face + status feedback", "oled", PURPLE), ("MAX98357A + speaker", "Digital audio output", "amp", AMBER), ("KY-040", "Volume + navigation", "encoder", BLUE), ("TTP223 ×2", "Happy / sad touch events", "touch", RED), ("ESP32-S3", "Edge controller + Wi-Fi", "esp32", GREEN)]
    for i,(title,desc,key,c) in enumerate(comps):
        x=0.7+(i%3)*4.15; y=1.65+(i//3)*2.55
        add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 3.65, 2.12, PANEL, "26364D")
        picture(s, assets[key], x+0.12, y+0.12, 1.42, 1.2, None)
        textbox(s, title, x+1.7, y+0.42, 1.75, 0.3, 12.5, c, True)
        textbox(s, desc, x+1.7, y+0.83, 1.65, 0.5, 10.7, MUTED)
    textbox(s, "Design principle: visible state + physical control + touch fallback + voice conversation.", 0.85, 6.55, 11.6, 0.25, 13, TEAL, True, align=PP_ALIGN.CENTER)

    # 7 touch interaction
    s = prs.slides.add_slide(blank); add_bg(s, "Interaction", 7, "Touch-driven affective interaction", "Two capacitive touch sensors create an immediate non-verbal response channel.")
    picture(s, assets["touch"], 0.8, 1.8, 3.5, 2.4, None)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 4.8, 1.8, 3.0, 1.15, PANEL, PURPLE)
    textbox(s, "TTP223 HAPPY", 5.05, 2.05, 2.5, 0.25, 14, PURPLE, True, align=PP_ALIGN.CENTER)
    textbox(s, "Tap → immediate reaction", 5.05, 2.4, 2.5, 0.22, 11, TEXT, align=PP_ALIGN.CENTER)
    add_arrow(s, 4.3, 3.0, 4.8, 2.35, PURPLE, 2)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 8.35, 1.8, 3.75, 3.95, PANEL_2, "2B4261")
    textbox(s, "Synchronized feedback", 8.65, 2.08, 3.1, 0.28, 14, TEAL, True, align=PP_ALIGN.CENTER)
    for i,(label,c) in enumerate([("OLED face", TEAL),("RGB color", GREEN),("short sound", AMBER),("state update", PURPLE)]):
        y=2.65+i*0.58
        add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 8.85, y, 2.75, 0.38, PANEL, c)
        textbox(s, label, 8.95, y+0.08, 2.55, 0.18, 10.7, TEXT, True, align=PP_ALIGN.CENTER)
    textbox(s, "No STT • no LLM • no TTS • no network wait", 4.35, 4.05, 3.95, 0.35, 13, AMBER, True, align=PP_ALIGN.CENTER)
    textbox(s, "The same state model keeps touch, voice, physical controls, display, LED, and audio consistent.", 1.0, 5.35, 10.9, 0.5, 18, TEXT, True, align=PP_ALIGN.CENTER)
    textbox(s, "System-level interaction contribution — not a new touch-sensing algorithm.", 1.0, 6.05, 10.9, 0.25, 10.5, MUTED, italic=True, align=PP_ALIGN.CENTER)

    # 8 software pipeline
    s = prs.slides.add_slide(blank); add_bg(s, "Software", 8, "Embedded + local AI software pipeline", "The heavier models remain replaceable while the device maintains deterministic feedback.")
    stages = [("1", "Capture", "INMP441\nI2S audio", TEAL), ("2", "Speech-to-text", "faster-whisper", PURPLE), ("3", "Reasoning", "Qwen3\nvia Ollama", AMBER), ("4", "Speech output", "Piper / Azure\nTTS", BLUE), ("5", "Playback", "MAX98357A\n+ speaker", GREEN)]
    for i,(n,title,body,c) in enumerate(stages):
        x=0.65+i*2.53
        add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.1, 2.0, 1.45, PANEL, c)
        textbox(s, n, x+0.12, 2.25, 0.3, 0.3, 16, c, True)
        textbox(s, title, x+0.45, 2.25, 1.4, 0.28, 12.5, TEXT, True)
        textbox(s, body, x+0.22, 2.75, 1.55, 0.55, 11.2, MUTED, align=PP_ALIGN.CENTER)
        if i < len(stages)-1: add_arrow(s, x+2.02, 2.82, x+2.45, 2.82, c, 1.8)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.85, 4.45, 5.7, 1.4, PANEL_2, TEAL)
    textbox(s, "Edge firmware", 1.15, 4.75, 1.55, 0.25, 13.5, TEAL, True)
    textbox(s, "GPIO events • I2C display • I2S transport • battery checks • LED state • server link", 1.15, 5.15, 4.9, 0.45, 11.1, TEXT)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 6.85, 4.45, 5.55, 1.4, PANEL_2, BLUE)
    textbox(s, "Server orchestration", 7.15, 4.75, 2.05, 0.25, 13.5, BLUE, True)
    textbox(s, "Session endpoint • timestamp logging • model substitution • failure recovery", 7.15, 5.15, 4.7, 0.45, 11.1, TEXT)

    # 9 simulation validation
    s = prs.slides.add_slide(blank); add_bg(s, "Validation", 9, "Wokwi simulation: what is verified", "The simulation is a functional-development aid, not a substitute for physical acoustic testing.")
    picture(s, assets["full"], 0.7, 1.65, 5.9, 4.95)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 7.0, 1.65, 5.35, 2.1, PANEL, TEAL)
    textbox(s, "Repeatable functional scenarios", 7.3, 1.92, 4.8, 0.28, 14, TEAL, True)
    rich_text(s, ["Idle startup and state recovery", "Happy touch / sad touch", "Encoder rotation and push", "Mode change and two-second sleep", "Battery reduction and response transitions"], 7.3, 2.35, 4.5, 1.05, 10.7)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 7.0, 4.05, 5.35, 2.05, PANEL, AMBER)
    textbox(s, "Boundary of the simulation", 7.3, 4.32, 4.7, 0.28, 14, AMBER, True)
    rich_text(s, ["Potentiometer approximates microphone activity", "Buzzer approximates audio playback", "Buttons approximate touch events", "INMP441 / MAX98357A / TTP223 behavior is not fully modeled"], 7.3, 4.75, 4.55, 1.05, 10.7)
    textbox(s, "Simulation evidence = GPIO, state transitions, display/LED behavior, event triggering.", 0.85, 6.62, 11.5, 0.22, 12.2, TEAL, True, align=PP_ALIGN.CENTER)

    # 10 evaluation
    s = prs.slides.add_slide(blank); add_bg(s, "Evaluation", 10, "Evaluation method and defensible results", "The paper separates functional correctness from performance measurements that require the physical prototype.")
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.75, 1.65, 5.75, 4.65, PANEL, GREEN)
    textbox(s, "Functional correctness", 1.05, 1.98, 4.8, 0.3, 15, GREEN, True, align=PP_ALIGN.CENTER)
    rich_text(s, ["Event detected", "Correct state transition", "OLED update", "RGB response", "Audio-event trigger", "Recovery after mode / sleep changes"], 1.25, 2.55, 4.6, 2.55, 13)
    textbox(s, "Binary pass / fail checks", 1.25, 5.65, 4.5, 0.25, 12, TEXT, True, align=PP_ALIGN.CENTER)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 6.85, 1.65, 5.75, 4.65, PANEL, BLUE)
    textbox(s, "Physical performance", 7.15, 1.98, 4.8, 0.3, 15, BLUE, True, align=PP_ALIGN.CENTER)
    rich_text(s, ["End-to-end response latency", "STT word error rate", "Command success rate", "False activation rate", "Battery current draw", "Sleep entry / exit time"], 7.35, 2.55, 4.6, 2.55, 13)
    textbox(s, "Measure on the prototype — do not infer from simulation", 7.25, 5.65, 4.9, 0.25, 12, TEXT, True, align=PP_ALIGN.CENTER)
    textbox(s, "Current result: a functional design and simulation model, not a completed acoustic benchmark.", 0.8, 6.62, 11.7, 0.22, 12.2, AMBER, True, align=PP_ALIGN.CENTER)

    # 11 discussion / limitations
    s = prs.slides.add_slide(blank); add_bg(s, "Discussion", 11, "Why the design matters — and where it stops", "LUNA provides two response classes: deterministic local events and conversational AI events.")
    card(s, 0.75, 1.7, 3.75, 2.05, "Human-centered", "Users can choose speech, touch, physical controls, or visual feedback depending on context and ability.", TEAL)
    card(s, 4.8, 1.7, 3.75, 2.05, "Reliability", "Simple actions still respond immediately when the AI pipeline is delayed or unavailable.", GREEN)
    card(s, 8.85, 1.7, 3.75, 2.05, "Privacy boundary", "Local STT/LLM can reduce external exposure, but authentication, logs, stored audio, and server access remain open questions.", PURPLE)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.75, 4.2, 11.85, 1.75, PANEL_2, "2B4261")
    textbox(s, "Known limitations", 1.0, 4.5, 2.1, 0.3, 14, AMBER, True)
    textbox(s, "Prototype design • no controlled user study • simulation does not model full acoustics/capacitance • MacBook-class host for practical real-time AI • Nepali TTS remains planned • touch novelty is integration-level", 1.0, 4.95, 10.9, 0.55, 12.3, TEXT)
    textbox(s, "The right claim is disciplined: LUNA demonstrates an implementable multimodal interaction architecture with a clear path to quantitative evaluation.", 1.0, 6.35, 10.9, 0.32, 13.2, TEAL, True, align=PP_ALIGN.CENTER)

    # 12 future / close
    s = prs.slides.add_slide(blank); add_bg(s, "Conclusion", 12, "Conclusion and next experimental steps", "A reproducible prototype architecture with a measurable research path.")
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.75, 1.7, 5.55, 3.95, PANEL, TEAL)
    textbox(s, "Takeaway", 1.05, 2.05, 4.85, 0.3, 16, TEAL, True, align=PP_ALIGN.CENTER)
    textbox(s, "LUNA combines voice, touch, physical controls, OLED feedback, RGB indication, and audio output around an ESP32-S3 — with local AI services kept modular.", 1.15, 2.65, 4.7, 1.25, 20, TEXT, True, align=PP_ALIGN.CENTER)
    textbox(s, "The central contribution is immediate affective touch response without invoking the speech pipeline.", 1.15, 4.35, 4.7, 0.6, 13, MUTED, align=PP_ALIGN.CENTER)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 6.75, 1.7, 5.75, 3.95, PANEL, PURPLE)
    textbox(s, "Next steps", 7.05, 2.05, 5.1, 0.3, 16, PURPLE, True, align=PP_ALIGN.CENTER)
    rich_text(s, ["Replace simulated audio/touch with physical modules", "Measure latency and STT error on a controlled test set", "Add wake-word and explicit privacy controls", "Extend touch vocabulary beyond happy / sad", "Evaluate English and Nepali interaction with users"], 7.25, 2.65, 4.75, 2.2, 12.5)
    textbox(s, "Thank you", 0.8, 6.15, 11.8, 0.45, 25, TEXT, True, align=PP_ALIGN.CENTER)
    textbox(s, "Questions / discussion", 0.8, 6.62, 11.8, 0.22, 11, TEAL, True, align=PP_ALIGN.CENTER)

    assert len(prs.slides) == 12
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    out = make_deck()
    print(out)

